from __future__ import annotations

import io
import os
import time
from typing import Optional, Tuple, Dict, Any, List

import httpx
from pypdf import PdfReader
from pptx import Presentation

# Optional advanced PDF + OCR stack
try:
    import pdfplumber  # type: ignore
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    from pdf2image import convert_from_bytes  # type: ignore
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False

try:
    import pytesseract  # type: ignore
    HAS_PYTESSERACT = True
except ImportError:
    HAS_PYTESSERACT = False


try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("⚠️  python-docx not installed. DOCX files will not be supported.")

    # Optional CLIP visual term extractor (scaffold)
    try:
        from models.clip_extractor import extract_visual_terms_from_pdf  # type: ignore
        HAS_CLIP_EXTRACTOR = True
    except Exception:
        HAS_CLIP_EXTRACTOR = False



async def download_file(file_url: str) -> Tuple[bytes, Optional[str]]:
    """Download a file and return bytes and detected content-type header if any."""
    async with httpx.AsyncClient(timeout=httpx.Timeout(60.0)) as client:
        response = await client.get(file_url)
        response.raise_for_status()
    return response.content, response.headers.get("content-type")


def extract_from_pdf(data: bytes) -> Tuple[str, Optional[int]]:
    """Extract text from PDF bytes using pypdf."""
    text_parts: list[str] = []
    pages = 0
    with io.BytesIO(data) as bio:
        reader = PdfReader(bio)
        pages = len(reader.pages)
        for page in reader.pages:
            try:
                text_parts.append(page.extract_text() or "")
            except Exception:
                # Continue on individual page errors
                continue
    text = "\n".join(filter(None, text_parts))
    
    # ADD VALIDATION
    if not text or to_word_count(text) < 10:
        text = "PDF_EXTRACTION_WARNING: Minimal text extracted. File may contain scanned images or be image-based.\n\n" + text
    
    return text, pages

def extract_from_pptx(data: bytes) -> str:
    """Extract text from PPTX bytes using python-pptx."""
    with io.BytesIO(data) as bio:
        prs = Presentation(bio)
        chunks: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    chunks.append(shape.text.strip())
        
        text = "\n".join(filter(None, chunks))
        
        # ADD VALIDATION
        if not text or to_word_count(text) < 5:
            text = "PPTX_EXTRACTION_WARNING: Minimal text found in presentation. Slides may be image-based.\n\n" + text
            
        return text

def extract_from_docx(data: bytes) -> str:
    """Extract text from DOCX files."""
    if not HAS_DOCX:
        return "DOCX_EXTRACTION_ERROR: python-docx library not installed. Please install: pip install python-docx"
    
    try:
        with io.BytesIO(data) as bio:
            doc = Document(bio)
            text_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            text = "\n".join(text_parts)
            
            # Validate extraction
            if not text or to_word_count(text) < 10:
                text = "DOCX_EXTRACTION_WARNING: Minimal text extracted from document.\n\n" + text
                
            return text
            
    except Exception as e:
        return f"DOCX_EXTRACTION_ERROR: {str(e)}"

def to_word_count(text: str) -> int:
    words = [w for w in text.replace("\n", " ").split(" ") if w.strip()]
    return len(words)


def _safe_set_tesseract_cmd() -> Optional[str]:
    """Configure pytesseract command path on Windows if provided via env.
    Returns the path used, or None if not set.
    """
    if not HAS_PYTESSERACT:
        return None
    tesseract_env = os.getenv("TESSERACT_CMD") or os.getenv("TESSERACT_PATH") or os.getenv("TESSERACT_EXE")
    if tesseract_env and os.path.isfile(tesseract_env):
        try:
            pytesseract.pytesseract.tesseract_cmd = tesseract_env
            return tesseract_env
        except Exception:
            return None
    return None


def _pdfplumber_extract(data: bytes) -> Tuple[str, int, Dict[str, Any]]:
    """Extract PDF text using pdfplumber with basic table stitching.
    Returns: (text, pages, metadata)
    metadata keys: method, per_page_word_counts
    """
    if not HAS_PDFPLUMBER:
        return "PDFPLUMBER_NOT_AVAILABLE", 0, {"method": "pdfplumber_unavailable", "per_page_word_counts": []}
    text_chunks: list[str] = []
    per_page_word_counts: list[int] = []
    per_page_texts: list[str] = []
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                # Basic table extraction: join table cells with pipe
                try:
                    tables = page.extract_tables() or []
                    for table in tables:
                        # table is list of rows; each row list of cells
                        row_lines = [" | ".join(cell for cell in row if cell) for row in table if any(cell and cell.strip() for cell in row)]
                        if row_lines:
                            page_text += "\n" + "\n".join(row_lines)
                except Exception:
                    pass
                if page_text.strip():
                    text_chunks.append(page_text.strip())
                per_page_texts.append(page_text.strip())
                per_page_word_counts.append(to_word_count(page_text))
        pages = len(per_page_word_counts)
        combined = "\n\n".join(text_chunks)
        return combined, pages, {"method": "pdfplumber", "per_page_word_counts": per_page_word_counts, "per_page_texts": per_page_texts}
    except Exception as e:
        return f"PDFPLUMBER_ERROR: {e}", 0, {"method": "pdfplumber_error", "per_page_word_counts": [], "per_page_texts": []}


def _batch_ranges(pages: List[int]) -> List[Tuple[int, int]]:
    """Group sorted page numbers into contiguous (start, end) ranges (1-based inclusive)."""
    if not pages:
        return []
    pages = sorted(set(pages))
    ranges: List[Tuple[int, int]] = []
    start = prev = pages[0]
    for p in pages[1:]:
        if p == prev + 1:
            prev = p
            continue
        ranges.append((start, prev))
        start = prev = p
    ranges.append((start, prev))
    return ranges


def _ocr_pdf_pages(
    data: bytes,
    pages_to_ocr: List[int],
    dpi: int = 150,
    poppler_path: Optional[str] = None,
    ocr_config: str = "--oem 3 --psm 6",
    batch_size: int = 20
) -> Tuple[Dict[int, str], int, float]:
    """OCR specific 1-based page numbers. Returns (page_texts, pages_processed, time_ms)."""
    if not (HAS_PDF2IMAGE and HAS_PYTESSERACT) or not pages_to_ocr:
        return {}, 0, 0.0
    start_t = time.time()
    page_texts: Dict[int, str] = {}
    processed = 0
    for rng in _batch_ranges(pages_to_ocr):
        rng_start, rng_end = rng
        # Batch further if very large range
        cur = rng_start
        while cur <= rng_end:
            last = min(cur + batch_size - 1, rng_end)
            try:
                images = convert_from_bytes(
                    data,
                    dpi=dpi,
                    first_page=cur,
                    last_page=last,
                    poppler_path=poppler_path
                )
            except Exception:
                break
            for idx, img in enumerate(images, start=cur):
                try:
                    txt = pytesseract.image_to_string(img, config=ocr_config) or ""
                except Exception:
                    txt = ""
                page_texts[idx] = (txt or "").strip()
                processed += 1
            cur = last + 1
    elapsed = (time.time() - start_t) * 1000.0
    return page_texts, processed, elapsed


def _normalize_text(text: str) -> str:
    """Post-process extracted text: normalize bullets, merge hyphenated wraps, collapse whitespace, dedupe."""
    if not text:
        return text
    # Normalize bullets
    text = (
        text.replace("\u2022", "- ")  # •
            .replace("•", "- ")
            .replace("●", "- ")
            .replace("–", "-")
            .replace("—", "-")
    )
    # Merge hyphenated line breaks: word-\nword -> wordword
    text = text.replace("-\n", "")
    # Collapse more than 2 blank lines
    while "\n\n\n" in text:
        text = text.replace("\n\n\n", "\n\n")
    # Remove simple consecutive duplicate lines
    lines = [l.rstrip() for l in text.splitlines()]
    dedup: List[str] = []
    prev = None
    for l in lines:
        if l != prev:
            dedup.append(l)
        prev = l
    return "\n".join(dedup)


async def extract_text_detailed(
    file_url: str,
    content_type_hint: Optional[str] = None,
    enable_ocr: bool = False,
    max_ocr_pages: int = 25,
    ocr_trigger_threshold: int = 10,
    ocr_word_gain_threshold: int = 20,
    poppler_path: Optional[str] = None,
    dpi: int = 150,
    ocr_config: str = "--oem 3 --psm 3",
    mode: str = "fast",
    full_ocr: bool = False,
    ocr_page_batch_size: int = 20,
    coverage_word_threshold: int = 10,
    coverage_target: float = 0.95
) -> Dict[str, Any]:
    """Enhanced extraction returning structured metadata.

    Returns dict with keys:
      text, pages, word_count, avg_words_per_page, extraction_method,
      ocr_triggered, ocr_pages, ocr_time_ms, warnings (list[str]), per_page_word_counts
    """
    data, detected = await download_file(file_url)
    mime = (content_type_hint or detected or "").lower()
    file_url_lower = file_url.lower()

    # Attempt to configure tesseract if path provided
    tesseract_used = _safe_set_tesseract_cmd()
    if not poppler_path:
        poppler_path = os.getenv("POPPLER_PATH")

    is_pdf = ("pdf" in mime) or file_url_lower.endswith(".pdf")
    is_pptx = any(file_url_lower.endswith(ext) for ext in (".ppt", ".pptx", ".pps", ".ppsx")) or "presentation" in mime
    is_docx = any(file_url_lower.endswith(ext) for ext in (".docx", ".doc")) or ("word" in mime or "document" in mime)

    warnings: list[str] = []
    metadata: Dict[str, Any] = {}
    text = ""
    pages: Optional[int] = None
    per_page_word_counts: list[int] = []
    per_page_texts: list[str] = []
    ocr_triggered = False
    ocr_pages = 0
    ocr_time_ms = 0.0
    extraction_method = "unknown"
    pages_ocrd_list: List[int] = []
    visual_terms: Dict[int, List[Tuple[str, float]]] = {}
    rescued_terms: List[str] = []

    try:
        if is_pdf:
            # First: pdfplumber path
            if HAS_PDFPLUMBER:
                pdf_text, pdf_pages, pdf_meta = _pdfplumber_extract(data)
                pages = pdf_pages
                per_page_word_counts = pdf_meta.get("per_page_word_counts", [])
                per_page_texts = pdf_meta.get("per_page_texts", [])
                base_wc = to_word_count(pdf_text)
                avg_pp = (sum(per_page_word_counts) / pages) if pages else 0.0
                extraction_method = pdf_meta.get("method", "pdfplumber")

                # Decide OCR trigger & perform selective OCR in 'complete' mode
                if (enable_ocr or mode == "complete") and pages:
                    # Determine pages to OCR
                    if full_ocr:
                        candidates = list(range(1, pages + 1))
                    else:
                        candidates = [i + 1 for i, wc in enumerate(per_page_word_counts) if wc < ocr_trigger_threshold]

                    if candidates:
                        # Respect max_ocr_pages if >0
                        if max_ocr_pages and max_ocr_pages > 0:
                            candidates = candidates[:max_ocr_pages]
                        pages_ocrd_list = list(candidates)
                        ocr_map, processed, elapsed = _ocr_pdf_pages(
                            data,
                            candidates,
                            dpi=dpi,
                            poppler_path=poppler_path,
                            ocr_config=ocr_config,
                            batch_size=ocr_page_batch_size,
                        )
                        ocr_pages = processed
                        ocr_time_ms = elapsed
                        if ocr_map:
                            ocr_triggered = True
                            # Merge OCR per page if it gains words
                            final_pages: List[str] = []
                            new_counts: List[int] = []
                            for idx in range(1, pages + 1):
                                raw_text = per_page_texts[idx - 1] if idx - 1 < len(per_page_texts) else ""
                                ocr_text = ocr_map.get(idx, "")
                                if ocr_text and (to_word_count(ocr_text) > to_word_count(raw_text) + ocr_word_gain_threshold or not raw_text.strip()):
                                    chosen = ocr_text
                                else:
                                    chosen = raw_text
                                final_pages.append(chosen)
                                new_counts.append(to_word_count(chosen))
                            per_page_texts = final_pages
                            per_page_word_counts = new_counts
                            text = "\n\n".join(p.strip() for p in final_pages if p and p.strip())
                            base_wc = to_word_count(text)
                            extraction_method = "pdfplumber+ocr-selective"
                        else:
                            text = pdf_text
                    else:
                        text = pdf_text
                else:
                    text = pdf_text

                word_count = base_wc
                avg_words_per_page = (word_count / pages) if pages else 0.0

                # CLIP-assisted visual term harvesting (Complete mode only)
                if mode == "complete" and pages and HAS_CLIP_EXTRACTOR:
                    try:
                        # Prefer low-density pages; if coverage below target, expand candidate set
                        candidates = [i + 1 for i, wc in enumerate(per_page_word_counts) if wc < ocr_trigger_threshold]
                        if not candidates:
                            candidates = list(range(1, min(pages, 5) + 1))
                        # If still below coverage target, widen search up to max_ocr_pages or 10
                        # coverage_pct computed later, so approximate using current low-density ratio
                        approx_covered = pages - len([wc for wc in per_page_word_counts if wc < coverage_word_threshold])
                        approx_cov_pct = (approx_covered / pages) if pages else 0.0
                        if approx_cov_pct < coverage_target:
                            widen = list(range(1, min(pages, max(10, max_ocr_pages or 10)) + 1))
                            # Merge keeping order and uniqueness
                            seen = set(candidates)
                            for p in widen:
                                if p not in seen:
                                    candidates.append(p)
                                    seen.add(p)
                        visual_terms = extract_visual_terms_from_pdf(
                            data,
                            pages=candidates,
                            dpi=max(dpi, 150),
                            poppler_path=poppler_path,
                        ) or {}
                        # Build rescued term list (not already in text)
                        text_lower = text.lower()
                        seen = set()
                        for _, pairs in visual_terms.items():
                            for label, score in pairs:
                                term = label.strip()
                                if not term:
                                    continue
                                if term.lower() not in text_lower and term.lower() not in seen:
                                    rescued_terms.append(term)
                                    seen.add(term.lower())
                        # Append small appendix so terms are present in downstream generation
                        if rescued_terms:
                            appendix = "\n\n## 📌 Visual Terms (images)\n" + "\n".join(f"- {t}" for t in rescued_terms[:30])
                            text = text + appendix
                            word_count = to_word_count(text)
                            base_wc = word_count
                    except Exception:
                        pass
            else:
                # Fallback pypdf
                pdf_text, pages = extract_from_pdf(data)
                text = pdf_text
                word_count = to_word_count(text)
                avg_words_per_page = (word_count / pages) if pages else 0.0
                extraction_method = "pypdf"
                warnings.append("pdfplumber_not_installed")
        elif is_pptx:
            text = extract_from_pptx(data)
            word_count = to_word_count(text)
            extraction_method = "pptx"
            avg_words_per_page = 0.0
        elif is_docx:
            text = extract_from_docx(data)
            word_count = to_word_count(text)
            extraction_method = "docx"
            avg_words_per_page = 0.0
        else:
            # Generic decode
            try:
                text = data.decode("utf-8")
            except Exception:
                text = data.decode(errors="ignore")
            word_count = to_word_count(text)
            extraction_method = "generic-bytes-decode"
            avg_words_per_page = 0.0
    except Exception as e:
        text = f"EXTRACTION_ERROR: {e}\n\n" + data.decode(errors="ignore")
        word_count = to_word_count(text)
        avg_words_per_page = 0.0
        warnings.append("exception_during_extraction")
        extraction_method = "error"

    # Post-process normalization
    text = _normalize_text(text)

    # Global validation & coverage
    if word_count < 20:
        warnings.append("low_total_word_count")
    if is_pdf and pages and pages > 0 and (sum(per_page_word_counts) if per_page_word_counts else word_count) / pages < ocr_trigger_threshold and not ocr_triggered:
        warnings.append("low_avg_words_per_page")
    if enable_ocr and not HAS_PYTESSERACT:
        warnings.append("pytesseract_not_installed")
    if enable_ocr and not HAS_PDF2IMAGE:
        warnings.append("pdf2image_not_installed")

    low_density_pages: List[int] = []
    coverage_pct: float = 0.0
    if is_pdf and pages and pages > 0 and per_page_word_counts:
        for i, wc in enumerate(per_page_word_counts, start=1):
            if wc < coverage_word_threshold:
                low_density_pages.append(i)
        covered = pages - len(low_density_pages)
        coverage_pct = round((covered / pages) * 100.0, 2)

    return {
        "text": text.strip(),
        "pages": pages,
        "word_count": word_count,
        "avg_words_per_page": avg_words_per_page,
        "extraction_method": extraction_method,
        "ocr_triggered": ocr_triggered,
        "ocr_pages": ocr_pages,
        "ocr_time_ms": ocr_time_ms,
        "warnings": warnings,
        "per_page_word_counts": per_page_word_counts,
        "tesseract_cmd_used": tesseract_used,
        "low_density_pages": low_density_pages,
        "coverage_pct": coverage_pct,
        "extraction_mode_used": f"{mode}|target={int(coverage_target*100)}%",
        "pages_ocrd": pages_ocrd_list,
        "visual_terms": visual_terms,
        "rescued_terms": rescued_terms,
        "rescued_terms_count": len(rescued_terms),
    }


async def extract_text_from_url(file_url: str, content_type_hint: Optional[str] = None) -> tuple[str, Optional[int], int]:
    """Backwards-compatible simple extraction wrapper (no OCR, minimal metadata)."""
    detailed = await extract_text_detailed(
        file_url=file_url,
        content_type_hint=content_type_hint,
        enable_ocr=False  # keep existing behavior
    )
    return detailed["text"], detailed["pages"], detailed["word_count"]