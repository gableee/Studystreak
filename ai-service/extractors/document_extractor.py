"""
Document text extraction utilities for PDF, PPTX, DOCX, and images (OCR).
Supports multiple extraction libraries for robustness.
"""

import logging
from typing import Optional, Dict, Any
from pathlib import Path
import io

logger = logging.getLogger(__name__)


class DocumentExtractor:
    """Extract text content from various document formats."""
    
    @staticmethod
    def extract_from_pdf(file_path_or_bytes, method: str = "pypdf") -> Dict[str, Any]:
        """
        Extract text from PDF file.
        
        Args:
            file_path_or_bytes: Path to PDF file or bytes object
            method: Extraction method - 'pypdf', 'pdfplumber', or 'pdfminer'
        
        Returns:
            dict with 'text', 'page_count', 'metadata'
        """
        try:
            if method == "pypdf":
                return DocumentExtractor._extract_pdf_pypdf(file_path_or_bytes)
            elif method == "pdfplumber":
                return DocumentExtractor._extract_pdf_pdfplumber(file_path_or_bytes)
            elif method == "pdfminer":
                return DocumentExtractor._extract_pdf_pdfminer(file_path_or_bytes)
            else:
                logger.warning(f"Unknown PDF method '{method}', defaulting to pypdf")
                return DocumentExtractor._extract_pdf_pypdf(file_path_or_bytes)
        except Exception as e:
            logger.error(f"PDF extraction failed with {method}: {e}")
            raise
    
    @staticmethod
    def _extract_pdf_pypdf(file_path_or_bytes) -> Dict[str, Any]:
        """Extract text using PyPDF2/pypdf."""
        try:
            from pypdf import PdfReader
            
            if isinstance(file_path_or_bytes, bytes):
                reader = PdfReader(io.BytesIO(file_path_or_bytes))
            else:
                reader = PdfReader(file_path_or_bytes)
            
            text_parts = []
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {page_num} ---\n{page_text}")
                except Exception as e:
                    logger.warning(f"Failed to extract page {page_num}: {e}")
            
            full_text = "\n\n".join(text_parts)
            
            # Extract metadata
            metadata = {}
            if reader.metadata:
                metadata = {
                    'title': reader.metadata.get('/Title', ''),
                    'author': reader.metadata.get('/Author', ''),
                    'subject': reader.metadata.get('/Subject', ''),
                    'creator': reader.metadata.get('/Creator', '')
                }
            
            logger.info(f"✅ Extracted {len(full_text)} chars from {len(reader.pages)} pages (pypdf)")
            
            return {
                'text': full_text,
                'page_count': len(reader.pages),
                'metadata': metadata,
                'method': 'pypdf'
            }
        except ImportError:
            logger.error("pypdf not installed. Install: pip install pypdf")
            raise
    
    @staticmethod
    def _extract_pdf_pdfplumber(file_path_or_bytes) -> Dict[str, Any]:
        """Extract text using pdfplumber (better for tables)."""
        try:
            import pdfplumber
            
            if isinstance(file_path_or_bytes, bytes):
                pdf = pdfplumber.open(io.BytesIO(file_path_or_bytes))
            else:
                pdf = pdfplumber.open(file_path_or_bytes)
            
            text_parts = []
            with pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"--- Page {page_num} ---\n{page_text}")
                    except Exception as e:
                        logger.warning(f"Failed to extract page {page_num}: {e}")
            
            full_text = "\n\n".join(text_parts)
            
            logger.info(f"✅ Extracted {len(full_text)} chars from {len(pdf.pages)} pages (pdfplumber)")
            
            return {
                'text': full_text,
                'page_count': len(pdf.pages),
                'metadata': {},
                'method': 'pdfplumber'
            }
        except ImportError:
            logger.error("pdfplumber not installed. Install: pip install pdfplumber")
            raise
    
    @staticmethod
    def _extract_pdf_pdfminer(file_path_or_bytes) -> Dict[str, Any]:
        """Extract text using pdfminer.six (robust for complex PDFs)."""
        try:
            from pdfminer.high_level import extract_text, extract_pages
            from pdfminer.pdfparser import PDFParser
            from pdfminer.pdfdocument import PDFDocument
            
            if isinstance(file_path_or_bytes, bytes):
                text = extract_text(io.BytesIO(file_path_or_bytes))
                # Count pages
                parser = PDFParser(io.BytesIO(file_path_or_bytes))
                document = PDFDocument(parser)
                page_count = sum(1 for _ in extract_pages(io.BytesIO(file_path_or_bytes)))
            else:
                text = extract_text(file_path_or_bytes)
                with open(file_path_or_bytes, 'rb') as f:
                    parser = PDFParser(f)
                    document = PDFDocument(parser)
                    page_count = sum(1 for _ in extract_pages(file_path_or_bytes))
            
            logger.info(f"✅ Extracted {len(text)} chars from {page_count} pages (pdfminer)")
            
            return {
                'text': text,
                'page_count': page_count,
                'metadata': {},
                'method': 'pdfminer'
            }
        except ImportError:
            logger.error("pdfminer.six not installed. Install: pip install pdfminer.six")
            raise
    
    @staticmethod
    def extract_from_pptx(file_path_or_bytes) -> Dict[str, Any]:
        """
        Extract text from PowerPoint (PPTX) file.
        
        Args:
            file_path_or_bytes: Path to PPTX file or bytes object
        
        Returns:
            dict with 'text', 'slide_count', 'metadata'
        """
        try:
            from pptx import Presentation
            
            if isinstance(file_path_or_bytes, bytes):
                prs = Presentation(io.BytesIO(file_path_or_bytes))
            else:
                prs = Presentation(file_path_or_bytes)
            
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
            
            full_text = "\n\n".join(text_parts)
            
            # Extract metadata
            metadata = {}
            if hasattr(prs.core_properties, 'title'):
                metadata['title'] = prs.core_properties.title or ''
            if hasattr(prs.core_properties, 'author'):
                metadata['author'] = prs.core_properties.author or ''
            
            logger.info(f"✅ Extracted {len(full_text)} chars from {len(prs.slides)} slides")
            
            return {
                'text': full_text,
                'slide_count': len(prs.slides),
                'metadata': metadata,
                'method': 'python-pptx'
            }
        except ImportError:
            logger.error("python-pptx not installed. Install: pip install python-pptx")
            raise
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise
    
    @staticmethod
    def extract_from_docx(file_path_or_bytes) -> Dict[str, Any]:
        """
        Extract text from Word (DOCX) file.
        
        Args:
            file_path_or_bytes: Path to DOCX file or bytes object
        
        Returns:
            dict with 'text', 'paragraph_count', 'metadata'
        """
        try:
            from docx import Document
            
            if isinstance(file_path_or_bytes, bytes):
                doc = Document(io.BytesIO(file_path_or_bytes))
            else:
                doc = Document(file_path_or_bytes)
            
            # Extract paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # Extract tables
            table_texts = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        table_texts.append(" | ".join(row_text))
            
            if table_texts:
                paragraphs.append("\n\n--- Tables ---\n" + "\n".join(table_texts))
            
            full_text = "\n\n".join(paragraphs)
            
            # Extract metadata
            metadata = {}
            if hasattr(doc.core_properties, 'title'):
                metadata['title'] = doc.core_properties.title or ''
            if hasattr(doc.core_properties, 'author'):
                metadata['author'] = doc.core_properties.author or ''
            if hasattr(doc.core_properties, 'subject'):
                metadata['subject'] = doc.core_properties.subject or ''
            
            logger.info(f"✅ Extracted {len(full_text)} chars from {len(doc.paragraphs)} paragraphs")
            
            return {
                'text': full_text,
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables),
                'metadata': metadata,
                'method': 'python-docx'
            }
        except ImportError:
            logger.error("python-docx not installed. Install: pip install python-docx")
            raise
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise
    
    @staticmethod
    def extract_from_image_ocr(file_path_or_bytes, language: str = 'eng') -> Dict[str, Any]:
        """
        Extract text from image using OCR (Tesseract).
        
        Args:
            file_path_or_bytes: Path to image file or bytes object
            language: OCR language (default: 'eng' for English)
        
        Returns:
            dict with 'text', 'confidence', 'metadata'
        """
        try:
            import pytesseract
            from PIL import Image
            
            if isinstance(file_path_or_bytes, bytes):
                image = Image.open(io.BytesIO(file_path_or_bytes))
            else:
                image = Image.open(file_path_or_bytes)
            
            # Extract text with detailed data for confidence scoring
            data = pytesseract.image_to_data(image, lang=language, output_type=pytesseract.Output.DICT)
            text = pytesseract.image_to_string(image, lang=language)
            
            # Calculate average confidence
            confidences = [int(conf) for conf in data['conf'] if conf != '-1']
            avg_confidence = sum(confidences) / len(confidences) if confidences else 0
            
            logger.info(f"✅ Extracted {len(text)} chars via OCR (confidence: {avg_confidence:.1f}%)")
            
            return {
                'text': text,
                'confidence': avg_confidence / 100.0,  # Normalize to 0-1
                'metadata': {
                    'language': language,
                    'image_size': f"{image.width}x{image.height}"
                },
                'method': 'tesseract-ocr'
            }
        except ImportError:
            logger.error("pytesseract or PIL not installed. Install: pip install pytesseract pillow")
            logger.error("Also ensure Tesseract OCR is installed on the system")
            raise
        except Exception as e:
            logger.error(f"OCR extraction failed: {e}")
            raise
    
    @staticmethod
    def extract_from_file(file_path_or_bytes, file_extension: Optional[str] = None) -> Dict[str, Any]:
        """
        Auto-detect file type and extract text.
        
        Args:
            file_path_or_bytes: Path to file or bytes object
            file_extension: Optional file extension hint (e.g., '.pdf', '.docx')
        
        Returns:
            dict with extracted text and metadata
        """
        try:
            # Detect file extension
            if file_extension is None:
                if isinstance(file_path_or_bytes, (str, Path)):
                    file_extension = Path(file_path_or_bytes).suffix.lower()
                else:
                    raise ValueError("file_extension required when passing bytes")
            
            file_extension = file_extension.lower()
            
            if file_extension in ['.pdf']:
                return DocumentExtractor.extract_from_pdf(file_path_or_bytes)
            elif file_extension in ['.pptx']:
                return DocumentExtractor.extract_from_pptx(file_path_or_bytes)
            elif file_extension in ['.docx']:
                return DocumentExtractor.extract_from_docx(file_path_or_bytes)
            elif file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
                return DocumentExtractor.extract_from_image_ocr(file_path_or_bytes)
            elif file_extension in ['.txt', '.md']:
                # Plain text files
                if isinstance(file_path_or_bytes, bytes):
                    text = file_path_or_bytes.decode('utf-8')
                else:
                    with open(file_path_or_bytes, 'r', encoding='utf-8') as f:
                        text = f.read()
                return {
                    'text': text,
                    'metadata': {},
                    'method': 'plain-text'
                }
            else:
                raise ValueError(f"Unsupported file extension: {file_extension}")
        
        except Exception as e:
            logger.error(f"File extraction failed: {e}")
            raise


def extract_text(file_path_or_bytes, file_extension: Optional[str] = None) -> str:
    """
    Convenience function to extract just the text content.
    
    Args:
        file_path_or_bytes: Path to file or bytes object
        file_extension: Optional file extension hint
    
    Returns:
        Extracted text as string
    """
    result = DocumentExtractor.extract_from_file(file_path_or_bytes, file_extension)
    return result.get('text', '')
